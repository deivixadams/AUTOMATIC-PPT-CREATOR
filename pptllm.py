from pptx import Presentation
from pptx.util import Inches
import os


def texto_a_presentacion(texto, slides):
    # Implementa la lógica para convertir el texto a una presentación.
    print("Texto a presentación")
    print("Texto:", texto)
    print("Slides:", slides)

def pdf_a_presentacion(texto, slides):
    # Implementa la lógica para convertir el texto del PDF a una presentación.
    print("PDF a presentación")
    print("Texto:", texto)
    print("Slides:", slides)


# Supón que info es la información procesada obtenida del LLM
info = [
    {"title": "Slide 1", "content": "Contenido del Slide 1", "image": "D:\AI\INFOTEP_3CASOS\PPT\PPT.LLM\imagen1.png"},
    {"title": "Slide 2", "content": "Contenido del Slide 2", "image": "D:\AI\INFOTEP_3CASOS\PPT\PPT.LLM\imagen2.png"},
]

# Crear una presentación
prs = Presentation()

for slide_info in info:
    # Añadir un slide con un título y contenido
    slide_layout = prs.slide_layouts[5]  # 5 es para un slide en blanco
    slide = prs.slides.add_slide(slide_layout)
    
    # Añadir título
    title_box = slide.shapes.title
    title_box.text = slide_info["title"]
    
    # Añadir contenido
    left = Inches(1)
    top = Inches(1.5)
    width = height = Inches(6)
    txBox=slide.shapes.add_textbox(left, top, width, height)
    tf=txBox.text_frame
    p=tf.add_paragraph()
    p.text=slide_info["content"]
    
    # Añadir imagen
    left = Inches(1)
    top = Inches(3)
    width = height = Inches(2)


    slide.shapes.add_picture(slide_info["image"], left, top, width, height)

# Guardar la presentación
prs.save(os.path.join("D:\AI\INFOTEP_3CASOS\PPT\PPT.LLM\RESULTADO", "mi_presentacion_con_imagenes.pptx"))
