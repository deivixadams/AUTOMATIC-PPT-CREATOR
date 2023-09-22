from pptx import Presentation
from pptx.util import Inches
import os
from datetime import datetime
import json

# Supón que info es la información procesada obtenida del LLM
info = [
    {"title": "Slide 1", "content": "Contenido del Slide 1", "image": "D:\AI\INFOTEP_3CASOS\PPT\PPT.LLM\imagen1.png"},
    {"title": "Slide 2", "content": "Contenido del Slide 2", "image": "D:\AI\INFOTEP_3CASOS\PPT\PPT.LLM\imagen2.png"},
]

def convertir_info(info):
    try:
        return json.loads(info)
    except json.JSONDecodeError:
        print("La cadena no está en un formato JSON válido.")


def crear_presentacion(info):
   infojson = convertir_info(info)
   print(f"Creando presentación... coño:\n")
   print(type(infojson))
   print("_"*50)
   print("\n"*2)
   prs = Presentation()
   ahora = datetime.now()
   fecha_hora = ahora.strftime("fecha-%d-%m-%Y_hora-%H-%M-%S")
   
   
   for slide_info in infojson:
        # Añadir un slide con un título y contenido
        slide_layout = prs.slide_layouts[5]  # 5 es para un slide en blanco
        slide = prs.slides.add_slide(slide_layout)
        
        # Añadir título
        title_box = slide.shapes.title
        title_box.text = slide_info["title"]
        
        # Añadir contenido
        left = Inches(1)
        top = Inches(1.5)
        width = height = Inches(6)
        txBox=slide.shapes.add_textbox(left, top, width, height)
        tf=txBox.text_frame
        p=tf.add_paragraph()
        p.text=slide_info["content"]

   nombre_archivo = f"{fecha_hora}.pptx"
   ruta_guardado = os.path.join("D:\AI\INFOTEP_3CASOS\PPT\PPT.LLM\RESULTADO", nombre_archivo)
   try:
      prs.save(ruta_guardado)
      if os.path.exists(ruta_guardado):
        print(f"He creado el archivo:{fecha_hora}")
   except PermissionError:
      print("No se pudo crear el archivo, probablemente esté abierto.") 

