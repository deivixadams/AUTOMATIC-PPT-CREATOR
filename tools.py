
import json
import os
import re
import nltk
from PyPDF2 import PdfReader # importamos la clase PdfReader
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime
import json
import openai
import tools
from langchain.chat_models import ChatOpenAI
from langchain import PromptTemplate
from langchain.llms import OpenAI
from langchain.chains import LLMChain, SimpleSequentialChain
import re
import requests
from bs4 import BeautifulSoup
import fitz  # fitz es el módulo principal de PyMuPDF



API_KEY = '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'
SERPER_KEY = 'e8644d7bc36c8d548e0a3e8813a8827d3083933f09fb742a784cbcce84e0f901'

class LimpiaTexto:
    def limpiar_texto(self, texto):
        pattern = re.compile('[^a-zA-Z0-9áéíóúÁÉÍÓÚñÑ.,?!¡¿ ]')
        texto_limpio = pattern.sub('', texto)
        return texto_limpio


class openai_langchain:
    
    def openai_lc(self):
        self.llm = ChatOpenAI(
            model_name="gpt-3.5-turbo",
            openai_api_key= API_KEY,
            temperature=1,
            max_tokens=900,
            model_kwargs={"top_p": 0.9, "frequency_penalty": 0.4}
        )
        return self.llm
    
    def completion(texto):
        completion = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=
            [
                {
                    "role" : "user", 
                    "content" : texto            
                }
            ]
        )
        
        return completion.choices[0].message["content"].strip()


class InfoConverter:
    def convertir_info(self, info):
        try:
            return json.loads(info)
        except json.JSONDecodeError:
            print("La cadena no está en un formato JSON válido.")

class PresentationCreator:

    def crear_presentacion(self, info):
        #infojson = self.convertir_info(info)
        print("Creando presentación...\n")
        prs = Presentation()
        ahora = datetime.now()
        fecha_hora = ahora.strftime("fecha-%d-%m-%Y_hora-%H-%M-%S")
        
        for slide_info in info:
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            title_box = slide.shapes.title
            title_box.text = slide_info["title"]
            left = Inches(1)
            top = Inches(1.5)
            width = height = Inches(6)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = slide_info["content"]
        
        nombre_archivo = f"{fecha_hora}.pptx"
        ruta_guardado = os.path.join("RESULTADO", nombre_archivo)
        try:
            prs.save(ruta_guardado)
            if os.path.exists(ruta_guardado):
                print(f"He creado el archivo:{fecha_hora} \n")
        except PermissionError:
            print("No se pudo crear el archivo, probablemente esté abierto.")



#procesando el pdf
class ExtractorTextoPDF:
    def __init__(self, ruta_del_pdf: str, paginas: str):
        self.ruta_del_pdf = ruta_del_pdf
        self.texto = ""
        self.paginas = self.parsear_paginas(paginas)
        self.extraer_texto()

        
    def parsear_paginas(self, paginas: str):
        # Convertir la cadena de paginas en una lista de enteros
        lista_paginas = []
        for x in paginas.split(','):
            if '-' in x:
                inicio, fin = map(int, x.split('-'))
                lista_paginas.extend(range(inicio, fin + 1))
            else:
                lista_paginas.append(int(x))
        return sorted(set(lista_paginas))
        
    def extraer_texto(self):
        try:
            doc = fitz.open(self.ruta_del_pdf)
            
            for num_pagina in self.paginas:
                if num_pagina > doc.page_count:
                    print(f"La página {num_pagina} no existe")
                    continue
                
                try:
                    pagina = doc.load_page(num_pagina - 1)  # -1 porque las páginas en PyMuPDF comienzan desde 0
                    self.texto += pagina.get_text()
                except Exception as e:
                    print(f"La página {num_pagina} no se puede extraer: {str(e)}")
                    
            doc.close()
            
        except Exception as e:
            print(f"Ocurrió un error al intentar abrir el PDF: {e}")

        #print(f"text PDF-------> {self.texto}")
        #input("Presione enter para continuar")
        return self.texto    

    '''
    def obtener_texto(self) -> str:
        print(f"text PDF-------> {self.texto}")
        return self.texto
    '''


#procesando la url
class WebReader:
    def read_web(self, url):  # Agregué un método para contener el bloque de código
        try:
            response = requests.get(url)
            if response.status_code == 200:
                soup = BeautifulSoup(response.text, 'html.parser')
                titulos = []
                for titulo in soup.find_all(['h1', 'h2']):
                    if len(titulos) < 1:
                        titulos.append(titulo.get_text())
                    else:
                        break
                
                titulo_limpio = re.sub('[^0-9a-zA-ZáéíóúÁÉÍÓÚñÑ\s-]+', '', titulos[0])
                titulo_limpio = titulo_limpio.strip()
                
                elementos_parrafo = [p for p in soup.find_all('p') if not p.find('a')]
                texto_completo = ''
                for elemento in elementos_parrafo:
                    texto_limpio = re.sub('[^0-9a-zA-ZáéíóúÁÉÍÓÚñÑ\s-]+', '', elemento.get_text())
                    texto_completo += texto_limpio + '\n' 
                    texto_completo = texto_completo.strip()[:7000]
                
                texto_final = """Del texto que te voy a enviar hacer un resumen con las conclusiones más importantes.
                solo utilizar los párrafos que formen oraciones con sentido completo.
                Debes responder en español.""" + titulo_limpio + " " + texto_completo
                
                # Aquí puedes continuar con tu lógica para llamar al API de OpenAI y procesar el resultado como prefieras.
                # Ejemplo (debes reemplazarlo por tu código real para interactuar con el API de OpenAI):
                # texto_ChatGPT = (completion(texto_final))
                
                #print("Resultado:\n", texto_final)  # Muestra el resultado por consola.
                return texto_final
            else:
                print("URL inválida, por favor introduzca otra URL")
        except Exception as e:
            print("Ha ocurrido un error al procesar la url:", e)