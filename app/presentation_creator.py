import os
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class PresentationCreator:
    def crear_presentacion(self, info, filepath):
        print("Creando presentación...\n")
        prs = Presentation()
        for slide_info in info:
            try:
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)

                # Agregar título
                title_box = slide.shapes.title
                title_box.text = slide_info["title"]
                title_box.text_frame.paragraphs[0].font.size = Pt(44)
                title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                # Agregar contenido
                left = Inches(1)
                top = Inches(2.5)
                width = Inches(8.5)
                height = Inches(3)
                content_box = slide.shapes.add_textbox(left, top, width, height)
                tf = content_box.text_frame
                tf.word_wrap = True  # Habilitar ajuste de palabras
                p = tf.add_paragraph()
                p.text = slide_info["content"]
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.LEFT  # Alinear el texto a la izquierda

                # Descargar y agregar imagen
                image_url = slide_info.get("image_url")
                if image_url:
                    try:
                        response = requests.get(image_url)
                        response.raise_for_status()  # Asegurarse de que la solicitud fue exitosa
                        image_stream = BytesIO(response.content)
                        img_left = Inches(1)
                        img_top = Inches(5)
                        img_width = Inches(8.5)
                        img_height = Inches(2)
                        slide.shapes.add_picture(image_stream, img_left, img_top, img_width, img_height)
                    except requests.RequestException as e:
                        print(f"Error al descargar la imagen: {e}")

            except Exception as e:
                print(f"Error al crear la diapositiva: {e}")
                raise e
        
        if not os.path.exists("RESULTADO"):
            os.makedirs("RESULTADO")
        try:
            prs.save(filepath)
            if os.path.exists(filepath):
                print(f"He creado el archivo: {filepath}\n")
        except PermissionError:
            print("No se pudo crear el archivo, probablemente esté abierto.")
        except Exception as e:
            print(f"Error al guardar la presentación: {e}")
            raise e
